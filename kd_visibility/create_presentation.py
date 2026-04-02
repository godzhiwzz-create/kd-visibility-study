#!/usr/bin/env python3
"""
KD Visibility论文演讲PPT生成器
生成15+页学术演讲PPT，多用图片和图表解释
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案
COLORS = {
    'primary': RGBColor(0x1E, 0x40, 0x7C),      # 深蓝
    'secondary': RGBColor(0x2E, 0x86, 0xAB),     # 青蓝
    'accent': RGBColor(0xF1, 0x8F, 0x01),        # 橙色强调
    'success': RGBColor(0x28, 0xA7, 0x45),       # 绿色
    'danger': RGBColor(0xDC, 0x35, 0x45),        # 红色
    'warning': RGBColor(0xFF, 0xC1, 0x07),       # 黄色
    'text': RGBColor(0x33, 0x33, 0x33),          # 深灰文字
    'light_bg': RGBColor(0xF8, 0xF9, 0xFA),      # 浅灰背景
}

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.3), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, section_num, title):
    """添加章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 左侧色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(4), prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    # 章节编号
    num_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.8), Inches(3), Inches(1.5)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}" if section_num < 10 else str(section_num)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 章节标题
    title_box = slide.shapes.add_textbox(
        Inches(4.5), Inches(3), Inches(8), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    return slide

def add_content_slide(prs, title, bullets, has_image=False, image_desc=""):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏背景
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 内容区域
    if has_image:
        # 左侧文字，右侧图片占位
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(6), Inches(5.5)
        )

        # 图片占位框
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.5),
            Inches(5.8), Inches(5.5)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = COLORS['light_bg']
        img_placeholder.line.color.rgb = COLORS['secondary']
        img_placeholder.line.width = Pt(2)

        # 图片说明
        if image_desc:
            desc_box = slide.shapes.add_textbox(
                Inches(7.2), Inches(6.5), Inches(5.4), Inches(0.8)
            )
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"📊 {image_desc}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['secondary']
            p.font.italic = True
    else:
        # 全宽文字
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5)
        )

    # 添加要点
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"● {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text']
        p.space_before = Pt(12)
        p.line_spacing = 1.5

    return slide

def add_figure_slide(prs, title, figure_desc, annotations):
    """添加图解页（大图+标注）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    # 主图区域
    main_fig = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3),
        Inches(8), Inches(5.5)
    )
    main_fig.fill.solid()
    main_fig.fill.fore_color.rgb = COLORS['light_bg']
    main_fig.line.color.rgb = COLORS['primary']
    main_fig.line.width = Pt(3)

    # 图中文字
    fig_text = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(7), Inches(1)
    )
    tf = fig_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"[ {figure_desc} ]"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['secondary']
    p.alignment = PP_ALIGN.CENTER

    # 右侧标注
    for i, (color, text) in enumerate(annotations):
        y_pos = 1.5 + i * 1.1

        # 颜色块
        color_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(9), Inches(y_pos),
            Inches(0.3), Inches(0.3)
        )
        color_box.fill.solid()
        color_box.fill.fore_color.rgb = color
        color_box.line.fill.background()

        # 标注文字
        anno_box = slide.shapes.add_textbox(
            Inches(9.5), Inches(y_pos), Inches(3.3), Inches(1)
        )
        tf = anno_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text']

    return slide

# ==================== 开始创建PPT ====================

print("正在生成PPT...")

# Slide 1: 封面
add_title_slide(prs,
    "可见度退化下知识蒸馏失效的分支机制分析",
    "Branch-wise Mechanism Analysis of Knowledge Distillation Failure\nUnder Visibility Degradation\n\n学术演讲 | 2026"
)

# Slide 2: 目录
slide = prs.slides.add_slide(prs.slide_layouts[6])
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
title.text_frame.paragraphs[0].text = "演讲大纲"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

outline_items = [
    ("01", "研究背景与动机", "为什么研究KD在雾天的失效？"),
    ("02", "研究框架设计", "如何设计实验来剖析问题？"),
    ("03", "核心发现", "Branch-wise性能结构"),
    ("04", "机制分析", "三个假设的验证"),
    ("05", "关键结论", "Occlusion是最直接支持的机制"),
    ("06", "启示与展望", "对未来研究的指导意义"),
]

for i, (num, title_text, desc) in enumerate(outline_items):
    y = 1.8 + i * 0.9

    # 编号
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(1), Inches(0.6))
    num_box.text_frame.paragraphs[0].text = num
    num_box.text_frame.paragraphs[0].font.size = Pt(28)
    num_box.text_frame.paragraphs[0].font.bold = True
    num_box.text_frame.paragraphs[0].font.color.rgb = COLORS['accent']

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1.8), Inches(y), Inches(4), Inches(0.6))
    title_box.text_frame.paragraphs[0].text = title_text
    title_box.text_frame.paragraphs[0].font.size = Pt(22)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = COLORS['text']

    # 描述
    desc_box = slide.shapes.add_textbox(Inches(1.8), Inches(y+0.35), Inches(8), Inches(0.5))
    desc_box.text_frame.paragraphs[0].text = desc
    desc_box.text_frame.paragraphs[0].font.size = Pt(16)
    desc_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# Slide 3: 章节分隔 - 背景
add_section_slide(prs, 1, "研究背景与动机")

# Slide 4: 背景 - 知识蒸馏简介
add_content_slide(prs,
    "知识蒸馏：模型压缩的核心技术",
    [
        "大型教师模型 → 小型学生模型：通过'软标签'传递知识",
        "在资源受限场景（边缘设备、实时检测）中至关重要",
        "Hinton 2015年提出，已成为工业界标准实践",
        "目标检测中用于压缩YOLO、Faster R-CNN等模型",
        "常见KD分支：Logit蒸馏、特征蒸馏、注意力蒸馏、定位蒸馏"
    ],
    has_image=True,
    image_desc="教师-学生框架示意图"
)

# Slide 5: 问题引入
add_content_slide(prs,
    "问题：雾天场景下KD失效了！",
    [
        "现状：KD在清晰图像上有效，但在雾天等退化条件下性能下降",
        "实际影响：自动驾驶、监控系统需要在恶劣天气下可靠运行",
        "核心困惑：为什么KD会失效？是 teacher 不行，还是 transfer 方式不对？",
        "研究空白：现有文献多关注'修复'，少有'机制剖析'",
        "本文目标：系统性地分析 KD 失效的根本原因"
    ],
    has_image=True,
    image_desc="清晰 vs 雾天检测效果对比"
)

# Slide 6: 章节分隔 - 框架
add_section_slide(prs, 2, "研究框架设计")

# Slide 7: 实验框架（大图）
slide = prs.slides.add_slide(prs.slide_layouts[6])
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
title.text_frame.paragraphs[0].text = "整体研究框架：观察 → 验证"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

# 框架图
framework = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3),
    Inches(12.3), Inches(5.8)
)
framework.fill.solid()
framework.fill.fore_color.rgb = COLORS['light_bg']
framework.line.color.rgb = COLORS['primary']
framework.line.width = Pt(2)

# Part 1 框
p1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(3.5), Inches(2))
p1.fill.solid()
p1.fill.fore_color.rgb = COLORS['secondary']
p1.line.fill.background()
p1_text = slide.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(3.1), Inches(1))
p1_text.text_frame.paragraphs[0].text = "Part 1\n结构化观察"
p1_text.text_frame.paragraphs[0].font.size = Pt(20)
p1_text.text_frame.paragraphs[0].font.bold = True
p1_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p1_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Part 2 框
p2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.8), Inches(3.5), Inches(2))
p2.fill.solid()
p2.fill.fore_color.rgb = COLORS['accent']
p2.line.fill.background()
p2_text = slide.shapes.add_textbox(Inches(8.2), Inches(2.3), Inches(3.1), Inches(1))
p2_text.text_frame.paragraphs[0].text = "Part 2\n机制验证"
p2_text.text_frame.paragraphs[0].font.size = Pt(20)
p2_text.text_frame.paragraphs[0].font.bold = True
p2_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p2_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 箭头
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5), Inches(2.5), Inches(2.5), Inches(0.6))
arrow.fill.solid()
arrow.fill.fore_color.rgb = COLORS['primary']
arrow.line.fill.background()

# 下方说明
flow_items = [
    ("5×3 实验矩阵", "5种KD分支 × 3种可见度"),
    ("机制假设", "M1分布不匹配 / M2遮挡 / M3不确定性"),
    ("统计验证", "相关性分析 + 因果干预"),
]

for i, (title_text, desc) in enumerate(flow_items):
    x = 1 + i * 4
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.3), Inches(3.8), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = COLORS['secondary']

    box_title = slide.shapes.add_textbox(Inches(x+0.1), Inches(4.4), Inches(3.6), Inches(0.5))
    box_title.text_frame.paragraphs[0].text = title_text
    box_title.text_frame.paragraphs[0].font.size = Pt(16)
    box_title.text_frame.paragraphs[0].font.bold = True
    box_title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    box_desc = slide.shapes.add_textbox(Inches(x+0.1), Inches(4.8), Inches(3.6), Inches(0.6))
    box_desc.text_frame.paragraphs[0].text = desc
    box_desc.text_frame.paragraphs[0].font.size = Pt(12)
    box_desc.text_frame.paragraphs[0].font.color.rgb = COLORS['text']

# Slide 8: 5×3 实验设计
add_content_slide(prs,
    "实验设计：5×3 矩阵",
    [
        "5种KD分支：Student-only / Logit / Feature / Attention / Localization",
        "3种可见度：Light (β=0.005) / Moderate (β=0.01) / Heavy (β=0.02)",
        "共15个训练实验，每个150 epochs，控制变量严格",
        "数据集：Foggy Cityscapes，8类目标检测",
        "教师：YOLOv8s，学生：YOLOv8n"
    ],
    has_image=True,
    image_desc="5×3实验矩阵热力图"
)

# Slide 9: 章节分隔 - 发现
add_section_slide(prs, 3, "核心发现")

# Slide 10: Branch-wise 发现（关键图）
slide = prs.slides.add_slide(prs.slide_layouts[6])
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
title.text_frame.paragraphs[0].text = "核心发现：Branch-wise 性能结构存在！"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

# 左侧文字
left = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(5))
tf = left.text_frame
tf.word_wrap = True
findings = [
    ("🔴 Logit KD", "最常用但效果最差，随可见度退化明显", COLORS['danger']),
    ("🟡 Feature KD", "中等表现，稳定性一般", COLORS['warning']),
    ("🟢 Localization KD", "相对最稳定，退化最缓慢", COLORS['success']),
    ("💡 启示", "不同分支对可见度敏感度不同！", COLORS['accent']),
]

for i, (label, text, color) in enumerate(findings):
    p = tf.add_paragraph()
    p.text = f"{label}: {text}"
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.space_before = Pt(12)

# 右侧图表区域
chart_area = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(1.2),
    Inches(7), Inches(5.5)
)
chart_area.fill.solid()
chart_area.fill.fore_color.rgb = COLORS['light_bg']
chart_area.line.color.rgb = COLORS['secondary']

chart_text = slide.shapes.add_textbox(Inches(7), Inches(3.5), Inches(4.6), Inches(1))
chart_text.text_frame.paragraphs[0].text = "[ 5×3性能曲线图 ]\n\nX轴: 可见度级别\nY轴: mAP@50\n\n展示Logit/Localization对比"
chart_text.text_frame.paragraphs[0].font.size = Pt(16)
chart_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
chart_text.text_frame.paragraphs[0].font.color.rgb = COLORS['secondary']

# Slide 11: 章节分隔 - 机制
add_section_slide(prs, 4, "机制分析")

# Slide 12: 三个假设
add_content_slide(prs,
    "三个竞争假设",
    [
        "M1 - 分布不匹配：雾天改变了特征分布，导致 teacher-student 对齐失效",
        "M2 - 遮挡：物体边界模糊导致定位困难，这是可见度退化的核心",
        "M3 - 不确定性：teacher 在雾天预测置信度下降，传递了错误信号",
        "验证方法：计算各机制指标与 KD Gain 的相关性",
        "结果：只有 M2 (遮挡) 显示出强相关性 (r=-0.989)"
    ],
    has_image=True,
    image_desc="三个机制的相关性分析图"
)

# Slide 13: 机制分析图（大图）
add_figure_slide(prs,
    "机制验证：只有Occlusion强相关！",
    "机制指标 vs KD Gain 散点图",
    [
        (COLORS['danger'], "M1 分布不匹配：|r|=0.3, p=0.7\n弱相关，无法解释"),
        (COLORS['warning'], "M3 不确定性：|r|=0.2, p=0.8\n同样无法解释"),
        (COLORS['success'], "M2 遮挡：r=-0.989, p=0.0015\n★ 强负相关！"),
    ]
)

# Slide 14: Causal Validation
add_content_slide(prs,
    "因果验证：干预实验",
    [
        "相关性≠因果性：需要干预实验确认 occlusion 是独立因素",
        "实验设计：固定 visibility (β=0)，人为控制 occlusion 比例",
        "Occlusion levels: 0, 0.1, 0.2, 0.3, 0.4, 0.5",
        "结果：KD Gain 随 occlusion 单调递减，确认因果关系",
        "结论：Occlusion 是 visibility degradation 中导致 KD 失效的独立因果因素"
    ],
    has_image=True,
    image_desc="Causal实验：Occlusion vs KD Gain曲线"
)

# Slide 15: 章节分隔 - 结论
add_section_slide(prs, 5, "关键结论")

# Slide 16: 主要结论
add_content_slide(prs,
    "主要研究结论",
    [
        "分支级结构确实存在：Localization KD 优于 Logit KD",
        "Occlusion 是最直接支持的失效机制 (r=-0.989, p=0.0015)",
        "M1/M3 (分布不匹配/不确定性) 解释力不足",
        "Logit KD 失效不是超参数问题，需要架构级解决方案",
        "定位信息 (Localization) 是雾天KD的关键"
    ],
    has_image=True,
    image_desc="结论框架图"
)

# Slide 17: 启示与展望
add_section_slide(prs, 6, "启示与展望")

# Slide 18: 启示
add_content_slide(prs,
    "对未来研究的启示",
    [
        "方法设计：雾天KD应优先处理遮挡和定位信息，而非仅分布对齐",
        "评估指标：应关注遮挡感知能力，而不仅是整体mAP",
        "教师选择：雾天场景下，高定位精度的teacher优于高分类精度的teacher",
        "扩展方向：其他退化类型（雨、雪、夜）的机制可能不同",
        "应用价值：为自动驾驶/监控系统在恶劣天气下的模型压缩提供指导"
    ],
    has_image=True,
    image_desc="未来研究方向"
)

# Slide 19: 感谢页
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 背景
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
    prs.slide_width, prs.slide_height
)
bg.fill.solid()
bg.fill.fore_color.rgb = COLORS['primary']
bg.line.fill.background()

# 感谢文字
thanks = slide.shapes.add_textbox(Inches(0), Inches(2.5), prs.slide_width, Inches(2))
tf = thanks.text_frame
tf.paragraphs[0].text = "感谢聆听"
tf.paragraphs[0].font.size = Pt(54)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 副标题
sub = slide.shapes.add_textbox(Inches(0), Inches(4.5), prs.slide_width, Inches(1))
tf = sub.text_frame
tf.paragraphs[0].text = "欢迎提问与讨论"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 论文信息
info = slide.shapes.add_textbox(Inches(0), Inches(6), prs.slide_width, Inches(0.8))
tf = info.text_frame
tf.paragraphs[0].text = "Branch-wise Mechanism Analysis of Knowledge Distillation Failure Under Visibility Degradation"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 保存
output_path = "/Users/godzhi/code/可见度视觉识别研究/kd_visibility/paper_spic/KD_Visibility_演讲.pptx"
prs.save(output_path)

print(f"PPT已生成: {output_path}")
print(f"总页数: {len(prs.slides)}")
print("\n包含页面:")
print("  1. 封面")
print("  2. 演讲大纲")
print("  3-5. 研究背景与动机")
print("  6-8. 研究框架设计")
print("  9-10. 核心发现")
print("  11-14. 机制分析")
print("  15-16. 关键结论")
print("  17-18. 启示与展望")
print("  19. 感谢页")
